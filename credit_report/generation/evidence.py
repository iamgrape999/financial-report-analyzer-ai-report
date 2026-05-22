from __future__ import annotations

import heapq
import logging
import time
from pathlib import Path
from typing import Optional

from credit_report.config import CREDIT_REPORTS_ROOT, CR_MAX_CHUNKS_PER_SECTION, SECTION_RETRIEVAL_KEYWORDS

logger = logging.getLogger(__name__)

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100

def _safe_report_dir(report_id: str) -> Path:
    """Resolve and validate report directory to prevent path traversal attacks.

    Uses a function-level import so that test patches of credit_report.config.CREDIT_REPORTS_ROOT
    are honoured (module-level binding would be frozen at import time).
    """
    from credit_report.config import CREDIT_REPORTS_ROOT as _root
    root = _root.resolve()
    candidate = (_root / report_id).resolve()
    if not str(candidate).startswith(str(root) + "/") and candidate != root:
        raise ValueError("Invalid report_id: path traversal attempt detected")
    return candidate


def _chunk_text(text: str) -> list[str]:
    """Split text into overlapping chunks, snapping to line boundaries where possible.

    Snapping to newlines prevents Markdown table rows from being split mid-row, which
    would cause Gemini to receive truncated column values without their header context.
    """
    if not text:
        return []
    chunks: list[str] = []
    i = 0
    while i < len(text):
        end = min(i + CHUNK_SIZE, len(text))
        # Snap to last newline within the overlap window to avoid mid-row splits
        if end < len(text):
            nl = text.rfind("\n", max(i + 1, end - CHUNK_OVERLAP), end)
            if nl > i:
                end = nl + 1
        chunk = text[i:end].strip()
        if chunk:
            chunks.append(chunk)
        # Advance by (chunk_length - overlap); if chunk is shorter than overlap,
        # advance to end to avoid infinite loop on very short texts.
        step = (end - i) - CHUNK_OVERLAP
        i += step if step > 0 else (end - i)
    return chunks


def _score_chunk(chunk: str, keywords: list[str]) -> int:
    """Count how many keywords appear in the chunk (case-insensitive)."""
    lower = chunk.lower()
    return sum(1 for kw in keywords if kw.lower() in lower)


def save_document_text(report_id: str, doc_id: str, text: str) -> None:
    """Persist extracted document text to CREDIT_REPORTS_ROOT/{report_id}/{doc_id}.txt."""
    doc_dir = CREDIT_REPORTS_ROOT / report_id
    doc_dir.mkdir(parents=True, exist_ok=True)
    (doc_dir / f"{doc_id}.txt").write_text(text, encoding="utf-8")


def save_document_binary(report_id: str, doc_id: str, file_bytes: bytes, filename: str) -> None:
    """Persist the original uploaded file bytes so ETL can be re-run without re-uploading."""
    doc_dir = CREDIT_REPORTS_ROOT / report_id
    doc_dir.mkdir(parents=True, exist_ok=True)
    (doc_dir / f"{doc_id}.bin").write_bytes(file_bytes)
    (doc_dir / f"{doc_id}.fname").write_text(filename, encoding="utf-8")


def load_document_texts(report_id: str) -> list[str]:
    """Load all extracted document texts for a report from the filesystem."""
    doc_dir = CREDIT_REPORTS_ROOT / report_id
    if not doc_dir.exists():
        return []
    texts: list[str] = []
    for txt_file in sorted(doc_dir.glob("*.txt")):
        try:
            texts.append(txt_file.read_text(encoding="utf-8"))
        except Exception as _e:
            logger.warning("load_document_texts: failed to read %s report=%s: %s", txt_file.name, report_id, _e)
    return texts


def retrieve_evidence(
    report_id: str,
    section_no: int,
    max_chunks: int = CR_MAX_CHUNKS_PER_SECTION,
) -> list[str]:
    """
    Return the most keyword-relevant text chunks for a section.

    Chunks are scored by how many of the section's retrieval keywords they contain,
    then the top-scoring chunks (up to max_chunks) are returned.
    """
    keywords = SECTION_RETRIEVAL_KEYWORDS.get(section_no, [])
    if not keywords:
        return []

    all_texts = load_document_texts(report_id)
    if not all_texts:
        return []

    # Bounded min-heap: keeps only top max_chunks scored chunks in memory.
    # (score, idx, chunk) — idx breaks score ties without comparing chunk strings.
    heap: list[tuple[int, int, str]] = []
    idx = 0
    for text in all_texts:
        for chunk in _chunk_text(text):
            score = _score_chunk(chunk, keywords)
            if score > 0:
                if len(heap) < max_chunks:
                    heapq.heappush(heap, (score, idx, chunk))
                elif score > heap[0][0]:
                    heapq.heapreplace(heap, (score, idx, chunk))
                idx += 1

    return [chunk for _, _, chunk in sorted(heap, key=lambda x: -x[0])]


# ── Text quality check ────────────────────────────────────────────────────────

def _text_quality_ok(text: str, min_meaningful_ratio: float = 0.05, min_chars: int = 80) -> bool:
    """Return True if text contains enough real content (letters, digits, CJK characters).

    Chinese brokerage PDFs often use CID fonts that pdfminer extracts as mostly whitespace
    or unmapped glyphs even though the byte count is large. This check prevents passing
    that garbage to ETL and ensures Vision OCR is used as fallback.
    """
    stripped = text.strip()
    if not stripped or len(stripped) < min_chars:
        return False
    meaningful = sum(
        1 for c in stripped
        if ('一' <= c <= '鿿')          # CJK Unified Ideographs
        or ('㐀' <= c <= '䶿')          # CJK Extension A
        or ('豈' <= c <= '﫿')          # CJK Compatibility Ideographs
        or c.isascii() and (c.isalpha() or c.isdigit())
    )
    ratio = meaningful / len(stripped)
    return ratio >= min_meaningful_ratio


def _quality_stats(text: str) -> dict:
    """Return quality metrics dict for logging."""
    stripped = text.strip()
    if not stripped:
        return {"chars": 0, "meaningful": 0, "ratio_pct": 0.0, "sample": ""}
    meaningful = sum(
        1 for c in stripped
        if ('一' <= c <= '鿿') or ('㐀' <= c <= '䶿') or ('豈' <= c <= '﫿')
        or c.isascii() and (c.isalpha() or c.isdigit())
    )
    return {
        "chars": len(stripped),
        "meaningful": meaningful,
        "ratio_pct": round(meaningful / len(stripped) * 100, 1),
        "sample": stripped[:120].replace("\n", " "),
    }


# ── PDF text extraction ───────────────────────────────────────────────────────

def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """Extract plain text from PDF bytes. Tries pdfminer first, then pypdf.

    Returns empty string if extraction fails or text quality is too low
    (e.g. CID-font Chinese PDFs where characters are unmapped).
    """
    pdf_kb = len(pdf_bytes) // 1024
    logger.info("[OCR] extract_text_from_pdf: start bytes=%dKB", pdf_kb)

    # pdfplumber — table-aware extraction, best for financial PDFs with structured tables
    t0 = time.perf_counter()
    plumber_text = _extract_tables_pdfplumber(pdf_bytes)
    elapsed = (time.perf_counter() - t0) * 1000
    if plumber_text and _text_quality_ok(plumber_text):
        logger.info("[OCR] pdfplumber: elapsed=%.0fms quality OK → using table-aware text", elapsed)
        return plumber_text
    if plumber_text:
        logger.info("[OCR] pdfplumber: elapsed=%.0fms low quality, trying pdfminer", elapsed)

    # pdfminer — best for native-text PDFs without complex table layouts
    t0 = time.perf_counter()
    try:
        import io
        from pdfminer.high_level import extract_text

        result = extract_text(io.BytesIO(pdf_bytes))
        elapsed = (time.perf_counter() - t0) * 1000
        stats = _quality_stats(result or "")
        logger.info(
            "[OCR] pdfminer: elapsed=%.0fms chars=%d meaningful=%d ratio=%.1f%% sample=%r",
            elapsed, stats["chars"], stats["meaningful"], stats["ratio_pct"], stats["sample"],
        )
        if result and _text_quality_ok(result):
            logger.info("[OCR] pdfminer: quality OK → using this text")
            return result
        if result and result.strip():
            logger.warning(
                "[OCR] pdfminer: LOW QUALITY text (ratio=%.1f%% < 5%% or chars<%d) "
                "— likely CID-font encoding issue, falling back to pypdf",
                stats["ratio_pct"], 80,
            )
        else:
            logger.info("[OCR] pdfminer: returned empty text")
    except Exception as e:
        logger.warning("[OCR] pdfminer: exception: %s", e)

    # pypdf fallback
    t0 = time.perf_counter()
    try:
        import io
        import pypdf

        reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
        page_count = len(reader.pages)
        pages = [page.extract_text() or "" for page in reader.pages]
        non_empty_pages = sum(1 for p in pages if p.strip())
        result = "\n\n".join(p for p in pages if p.strip())
        elapsed = (time.perf_counter() - t0) * 1000
        stats = _quality_stats(result)
        logger.info(
            "[OCR] pypdf: elapsed=%.0fms pages=%d non_empty_pages=%d chars=%d "
            "meaningful=%d ratio=%.1f%% sample=%r",
            elapsed, page_count, non_empty_pages,
            stats["chars"], stats["meaningful"], stats["ratio_pct"], stats["sample"],
        )
        if result and _text_quality_ok(result):
            logger.info("[OCR] pypdf: quality OK → using this text")
            return result
        if result and result.strip():
            logger.warning(
                "[OCR] pypdf: LOW QUALITY text (ratio=%.1f%%) — will try Vision OCR",
                stats["ratio_pct"],
            )
        else:
            logger.info("[OCR] pypdf: returned empty text")
    except Exception as e:
        logger.warning("[OCR] pypdf: exception: %s", e)

    logger.warning(
        "[OCR] extract_text_from_pdf: all parsers failed or produced low-quality text "
        "— returning empty so caller triggers Vision OCR"
    )
    return ""


_PDF_SPARSE_PAGE_THRESHOLD = 50   # chars below which a page is considered "image-only"
_PDF_MAX_IMAGE_PAGE_OCR = 5       # max Vision OCR calls per PDF for sparse pages


def _extract_tables_pdfplumber(pdf_bytes: bytes) -> str:
    """Extract text + tables from PDF via pdfplumber, rendering tables as Markdown.

    Pages with embedded images but very little native text (image-only pages, e.g.
    financial tables rendered as pictures) are sent to Gemini Vision OCR so that
    embedded financial data is not silently lost.
    """
    try:
        import io
        import pdfplumber

        parts: list[str] = []
        image_ocr_count = 0
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page_no, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        if not table or not table[0]:
                            continue
                        headers = [str(c).strip() if c is not None else "" for c in table[0]]
                        col_context = " | ".join(h for h in headers if h)
                        # Structured table opener — survives chunk splits downstream
                        parts.append(f"\n[TABLE page={page_no} columns={col_context}]\n")
                        parts.append("| " + " | ".join(headers) + " |")
                        parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
                        for row_idx, row in enumerate(table[1:]):
                            # Repeat column context every 15 data rows so that even when
                            # Gemini receives a mid-table chunk, it knows the column mapping.
                            if row_idx > 0 and row_idx % 15 == 0:
                                parts.append(f"<!-- columns: {col_context} -->")
                            cells = [str(c).strip() if c is not None else "" for c in row]
                            parts.append("| " + " | ".join(cells) + " |")
                if page_text.strip():
                    parts.append(f"\n[Page {page_no}]\n{page_text}")
                elif (
                    image_ocr_count < _PDF_MAX_IMAGE_PAGE_OCR
                    and len(page_text) < _PDF_SPARSE_PAGE_THRESHOLD
                    and page.images  # pdfplumber detected embedded images on this page
                ):
                    # Image-only page: send to Gemini Vision so charts/tables aren't lost
                    try:
                        page_pdf_bytes = _render_single_pdf_page(pdf_bytes, page_no - 1)
                        if page_pdf_bytes:
                            vis_text = extract_text_from_image(page_pdf_bytes, "application/pdf")
                            if vis_text.strip():
                                parts.append(f"\n[Page {page_no} Vision OCR]\n{vis_text}")
                                image_ocr_count += 1
                                logger.info("[OCR] pdfplumber: page %d image-only → Vision OCR (%d chars)", page_no, len(vis_text))
                    except Exception as page_ocr_e:
                        logger.warning("[OCR] pdfplumber: page %d Vision OCR failed: %s", page_no, page_ocr_e)

        result = "\n".join(parts)
        logger.info("[OCR] pdfplumber: extracted %d chars from PDF (image_pages_ocr=%d)", len(result), image_ocr_count)
        return result
    except ImportError:
        return ""
    except Exception as e:
        logger.warning("[OCR] pdfplumber: extraction failed: %s", e)
        return ""


def _render_single_pdf_page(pdf_bytes: bytes, page_idx: int) -> bytes:
    """Extract a single PDF page as a minimal PDF blob for Vision OCR."""
    try:
        import io
        from pypdf import PdfReader, PdfWriter

        reader = PdfReader(io.BytesIO(pdf_bytes))
        if page_idx >= len(reader.pages):
            return b""
        writer = PdfWriter()
        writer.add_page(reader.pages[page_idx])
        buf = io.BytesIO()
        writer.write(buf)
        return buf.getvalue()
    except Exception as e:
        logger.warning("[OCR] _render_single_pdf_page failed: %s", e)
        return b""


# ── Office format extraction ──────────────────────────────────────────────────

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract plain text + embedded images from DOCX bytes.

    Paragraphs and native tables are extracted via python-docx.
    Embedded images (charts rendered as pictures, scanned exhibits) are
    sent to Gemini Vision OCR so image-based financial data is captured.
    """
    logger.info("[OCR] extract_text_from_docx: start bytes=%dKB", len(file_bytes) // 1024)
    try:
        import io
        from docx import Document

        _DOCX_MAX_IMAGE_OCR = 10
        _DOCX_MIN_IMAGE_BYTES = 10 * 1024   # skip tiny images (icons, logos)

        t0 = time.perf_counter()
        doc = Document(io.BytesIO(file_bytes))
        parts: list[str] = []
        para_count = 0
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
                para_count += 1
        table_rows = 0
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
                    table_rows += 1

        # Extract embedded images and run Vision OCR on each
        images_ocr_count = 0
        try:
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            seen_rids: set[str] = set()
            for rel in doc.part.rels.values():
                if images_ocr_count >= _DOCX_MAX_IMAGE_OCR:
                    break
                if rel.reltype != RT.IMAGE:
                    continue
                rid = rel.rId
                if rid in seen_rids:
                    continue
                seen_rids.add(rid)
                try:
                    image_bytes = rel.target_part.blob
                    if len(image_bytes) < _DOCX_MIN_IMAGE_BYTES:
                        continue
                    mt = rel.target_part.content_type or "image/png"
                    if mt not in ("image/jpeg", "image/png", "image/gif", "image/webp"):
                        mt = "image/png"
                    img_text = extract_text_from_image(image_bytes, mt)
                    if img_text.strip():
                        parts.append(f"\n[Embedded Image {images_ocr_count + 1}]\n{img_text}")
                        images_ocr_count += 1
                except Exception as img_e:
                    logger.warning("[OCR] docx embedded image OCR failed: %s", img_e)
        except Exception as img_loop_e:
            logger.warning("[OCR] docx image loop failed: %s", img_loop_e)

        result = "\n".join(parts)
        elapsed = (time.perf_counter() - t0) * 1000
        stats = _quality_stats(result)
        logger.info(
            "[OCR] extract_text_from_docx: elapsed=%.0fms paragraphs=%d table_rows=%d "
            "images_ocr=%d chars=%d ratio=%.1f%% sample=%r",
            elapsed, para_count, table_rows, images_ocr_count,
            stats["chars"], stats["ratio_pct"], stats["sample"],
        )
        return result
    except Exception as e:
        logger.warning("[OCR] extract_text_from_docx: failed: %s", e)
        return ""


_PPTX_MAX_IMAGE_OCR = 20      # max Vision OCR calls per PPTX file
_PPTX_MIN_IMAGE_BYTES = 10 * 1024  # skip tiny images < 10 KB (icons/logos)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract plain text from PPTX bytes using python-pptx.

    Text shapes and tables are extracted via python-pptx.  Image shapes
    (charts rendered as pictures, scanned slides, infographic slides) are
    sent to Gemini Vision OCR so embedded financial data is not lost.
    """
    logger.info("[OCR] extract_text_from_pptx: start bytes=%dKB", len(file_bytes) // 1024)
    try:
        import io
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        t0 = time.perf_counter()
        prs = Presentation(io.BytesIO(file_bytes))
        parts: list[str] = []
        slide_count = len(prs.slides)
        shape_count = 0
        images_ocr_count = 0
        for slide_no, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {slide_no}]"]
            for shape in slide.shapes:
                shape_count += 1
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)
                # Vision OCR for embedded image shapes (charts, scanned slides, infographics)
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and images_ocr_count < _PPTX_MAX_IMAGE_OCR
                ):
                    try:
                        img_blob = shape.image.blob
                        if len(img_blob) >= _PPTX_MIN_IMAGE_BYTES:
                            img_mime = shape.image.content_type or "image/jpeg"
                            img_text = extract_text_from_image(img_blob, img_mime)
                            if img_text.strip():
                                slide_parts.append(f"[Image OCR]\n{img_text.strip()}")
                                images_ocr_count += 1
                    except Exception as _img_err:
                        logger.debug(
                            "[OCR] pptx: image OCR failed slide=%d shape=%r: %s",
                            slide_no, getattr(shape, "name", "?"), _img_err,
                        )
            if len(slide_parts) > 1:
                parts.append("\n".join(slide_parts))
        result = "\n\n".join(parts)
        elapsed = (time.perf_counter() - t0) * 1000
        stats = _quality_stats(result)
        logger.info(
            "[OCR] extract_text_from_pptx: elapsed=%.0fms slides=%d shapes=%d "
            "non_empty_slides=%d images_ocr=%d chars=%d ratio=%.1f%% sample=%r",
            elapsed, slide_count, shape_count, len(parts), images_ocr_count,
            stats["chars"], stats["ratio_pct"], stats["sample"],
        )
        return result
    except Exception as e:
        logger.warning("[OCR] extract_text_from_pptx: failed: %s", e)
        return ""


# ── Vision OCR ────────────────────────────────────────────────────────────────

_OCR_PROMPT = (
    "Extract ALL text, numbers, tables, and structured data from this document. "
    "The document may be in Traditional Chinese (繁體中文), Simplified Chinese (简体中文), "
    "or English — extract all text faithfully in its original language. "
    "For tables, preserve the structure using | separators between columns and "
    "new lines between rows. Include all financial figures, dates, company names, "
    "key metrics, percentages, ratios, and headers exactly as shown. "
    "Preserve original currency symbols and units (NTD, TWD, USD, HKD, etc.). "
    "Output as plain text maintaining the original layout and page structure."
)


def extract_text_from_image(image_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    """Use Gemini Vision to OCR and extract text+tables from an image."""
    logger.info("[OCR] extract_text_from_image: start bytes=%dKB mime=%s", len(image_bytes) // 1024, mime_type)
    try:
        from google import genai
        from google.genai import types as genai_types
        from credit_report.config import GEMINI_API_KEY, GEMINI_OCR_MODEL

        if not GEMINI_API_KEY:
            logger.warning("[OCR] extract_text_from_image: GEMINI_API_KEY not set — cannot OCR")
            return ""

        t0 = time.perf_counter()
        client = genai.Client(api_key=GEMINI_API_KEY)
        logger.info("[OCR] extract_text_from_image: calling Gemini model=%s max_tokens=32768", GEMINI_OCR_MODEL)
        response = client.models.generate_content(
            model=GEMINI_OCR_MODEL,
            contents=[
                genai_types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
                genai_types.Part.from_text(_OCR_PROMPT),
            ],
            config=genai_types.GenerateContentConfig(max_output_tokens=32768, thinking_config=genai_types.ThinkingConfig(thinking_budget=0)),
        )
        elapsed = (time.perf_counter() - t0) * 1000
        result = response.text or ""
        finish_reason = getattr(getattr(response, "candidates", [None])[0], "finish_reason", "unknown") if response.candidates else "no_candidates"
        stats = _quality_stats(result)
        logger.info(
            "[OCR] extract_text_from_image: elapsed=%.0fms chars=%d meaningful=%d "
            "ratio=%.1f%% finish_reason=%s sample=%r",
            elapsed, stats["chars"], stats["meaningful"],
            stats["ratio_pct"], finish_reason, stats["sample"],
        )
        if not result:
            logger.warning("[OCR] extract_text_from_image: Gemini returned empty response")
        return result
    except Exception as e:
        logger.warning("[OCR] extract_text_from_image: Gemini Vision failed: %s", e)
        return ""


def _split_pdf_pages(pdf_bytes: bytes) -> list[bytes]:
    """Return a list of single-page PDF byte blobs using pypdf."""
    try:
        import io
        from pypdf import PdfReader, PdfWriter

        reader = PdfReader(io.BytesIO(pdf_bytes))
        pages: list[bytes] = []
        for page in reader.pages:
            writer = PdfWriter()
            writer.add_page(page)
            buf = io.BytesIO()
            writer.write(buf)
            pages.append(buf.getvalue())
        return pages
    except Exception as e:
        logger.warning("[OCR] _split_pdf_pages failed: %s — will send full PDF", e)
        return []


def _iter_pdf_pages(pdf_bytes: bytes):
    """Yield single-page PDF byte blobs one at a time (memory-efficient alternative to _split_pdf_pages)."""
    try:
        import io
        from pypdf import PdfReader, PdfWriter

        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            writer = PdfWriter()
            writer.add_page(page)
            buf = io.BytesIO()
            writer.write(buf)
            yield buf.getvalue()
    except Exception as e:
        logger.warning("[OCR] _iter_pdf_pages failed: %s — stopping iteration", e)
        return


# Gemini inline PDF upload limit per request
_GEMINI_PDF_SIZE_LIMIT = 20 * 1024 * 1024  # 20 MB
# Page count threshold: PDFs larger than this go through page-by-page VLM
_PDF_PAGE_THRESHOLD = 30


def extract_text_from_scanned_pdf_vision(pdf_bytes: bytes, max_pages: int = 200) -> str:
    """
    For scanned/image PDFs where text extraction yields nothing:
    send the PDF to Gemini Vision for OCR.

    Strategy:
    - Small PDFs (≤ 20 MB): send whole file in one call.
    - Large PDFs (> 20 MB) or those with many pages: split into single-page chunks,
      OCR each page separately, concatenate.  This prevents silent truncation and
      ensures every page is processed.
    """
    pdf_kb = len(pdf_bytes) // 1024
    logger.info("[OCR] extract_text_from_scanned_pdf_vision: start pdf_kb=%d", pdf_kb)
    try:
        from google import genai
        from google.genai import types as genai_types
        from credit_report.config import GEMINI_API_KEY, GEMINI_OCR_MODEL

        if not GEMINI_API_KEY:
            logger.warning("[OCR] extract_text_from_scanned_pdf_vision: GEMINI_API_KEY not set")
            return ""

        client = genai.Client(api_key=GEMINI_API_KEY)
        t0 = time.perf_counter()

        # Decide strategy: whole-PDF vs page-by-page
        use_page_split = len(pdf_bytes) > _GEMINI_PDF_SIZE_LIMIT
        if not use_page_split:
            # Try to count pages; switch to page-split if > _PDF_PAGE_THRESHOLD
            try:
                import io
                from pypdf import PdfReader
                n_pages = len(PdfReader(io.BytesIO(pdf_bytes)).pages)
                use_page_split = n_pages > _PDF_PAGE_THRESHOLD
                logger.info("[OCR] PDF pages=%d use_page_split=%s", n_pages, use_page_split)
            except Exception:
                pass

        if not use_page_split:
            # Single-call path (fast for short PDFs)
            logger.info("[OCR] calling Gemini Vision PDF OCR (single-call) model=%s", GEMINI_OCR_MODEL)
            response = client.models.generate_content(
                model=GEMINI_OCR_MODEL,
                contents=[
                    genai_types.Part.from_bytes(data=pdf_bytes, mime_type="application/pdf"),
                    genai_types.Part.from_text(_OCR_PROMPT),
                ],
                config=genai_types.GenerateContentConfig(max_output_tokens=32768, thinking_config=genai_types.ThinkingConfig(thinking_budget=0)),
            )
            result = response.text or ""
            elapsed = (time.perf_counter() - t0) * 1000
            stats = _quality_stats(result)
            finish_reason = (
                str(response.candidates[0].finish_reason) if response.candidates else "no_candidates"
            )
            logger.info(
                "[OCR] vision_ocr single-call: elapsed=%.0fms chars=%d "
                "meaningful=%d ratio=%.1f%% finish_reason=%s sample=%r",
                elapsed, stats["chars"], stats["meaningful"],
                stats["ratio_pct"], finish_reason, stats["sample"][:200],
            )
            if stats["ratio_pct"] < 5:
                logger.warning(
                    "[OCR] vision_ocr: low quality ratio=%.1f%% — escalating to page-by-page",
                    stats["ratio_pct"],
                )
                use_page_split = True  # fall through to page split below
            else:
                return result

        # Page-by-page path: iterate pages one at a time to avoid holding all blobs in memory.
        # Count pages first with PdfReader (no blob allocation), then stream via _iter_pdf_pages.
        try:
            import io as _io2
            from pypdf import PdfReader as _PR2
            n_pages = len(_PR2(_io2.BytesIO(pdf_bytes)).pages)
        except Exception:
            n_pages = 0

        if n_pages == 0:
            # pypdf can't read the file — fall back to capped single-call
            data = pdf_bytes[:_GEMINI_PDF_SIZE_LIMIT]
            logger.info("[OCR] page split failed — fallback single-call sending_kb=%d", len(data) // 1024)
            response = client.models.generate_content(
                model=GEMINI_OCR_MODEL,
                contents=[
                    genai_types.Part.from_bytes(data=data, mime_type="application/pdf"),
                    genai_types.Part.from_text(_OCR_PROMPT),
                ],
                config=genai_types.GenerateContentConfig(max_output_tokens=32768, thinking_config=genai_types.ThinkingConfig(thinking_budget=0)),
            )
            return response.text or ""

        pages_to_process = min(n_pages, max_pages)
        logger.info(
            "[OCR] page-by-page OCR: total_pages=%d processing=%d model=%s",
            n_pages, pages_to_process, GEMINI_OCR_MODEL,
        )
        page_texts: list[str] = []
        for page_idx, page_blob in enumerate(_iter_pdf_pages(pdf_bytes)):
            if page_idx >= max_pages:
                break
            try:
                page_response = client.models.generate_content(
                    model=GEMINI_OCR_MODEL,
                    contents=[
                        genai_types.Part.from_bytes(data=page_blob, mime_type="application/pdf"),
                        genai_types.Part.from_text(_OCR_PROMPT),
                    ],
                    config=genai_types.GenerateContentConfig(max_output_tokens=8192, thinking_config=genai_types.ThinkingConfig(thinking_budget=0)),
                )
                page_text = page_response.text or ""
                page_texts.append(page_text)
                if page_idx % 10 == 0:
                    logger.info("[OCR] page-by-page: processed %d/%d pages", page_idx + 1, pages_to_process)
            except Exception as page_err:
                logger.warning("[OCR] page-by-page: page %d failed: %s", page_idx + 1, page_err)
                page_texts.append("")  # preserve page numbering

        result = "\n\n".join(t for t in page_texts if t.strip())
        elapsed = (time.perf_counter() - t0) * 1000
        stats = _quality_stats(result)
        logger.info(
            "[OCR] page-by-page OCR complete: total_pages=%d elapsed=%.0fms chars=%d "
            "meaningful=%d ratio=%.1f%% sample=%r",
            pages_to_process, elapsed, stats["chars"], stats["meaningful"],
            stats["ratio_pct"], stats["sample"][:200],
        )
        if n_pages > max_pages:
            logger.warning(
                "[OCR] page-by-page: PDF has %d pages but only first %d were processed",
                n_pages, max_pages,
            )
        return result

    except Exception as e:
        logger.warning("[OCR] extract_text_from_scanned_pdf_vision: Gemini Vision failed: %s", e)
        return ""


async def extract_text_from_scanned_pdf_vision_async(
    pdf_bytes: bytes,
    on_progress: Optional[callable] = None,
    max_pages: int = 200,
) -> str:
    """
    Async page-by-page Gemini Vision OCR for large PDFs (up to 200 pages).
    on_progress(page_idx, total_pages, chars_so_far) called after each page.
    """
    from google import genai
    from google.genai import types as genai_types
    from credit_report.config import GEMINI_API_KEY, GEMINI_OCR_MODEL

    if not GEMINI_API_KEY:
        return ""

    pdf_kb = len(pdf_bytes) // 1024
    logger.info("[OCR-ASYNC] start pdf_kb=%d max_pages=%d", pdf_kb, max_pages)
    t0 = time.perf_counter()

    # For small PDFs that fit in one call, use single-call path
    use_page_split = len(pdf_bytes) > _GEMINI_PDF_SIZE_LIMIT
    if not use_page_split:
        try:
            import io
            from pypdf import PdfReader
            n_pages = len(PdfReader(io.BytesIO(pdf_bytes)).pages)
            use_page_split = n_pages > _PDF_PAGE_THRESHOLD
            logger.info("[OCR-ASYNC] pages=%d use_page_split=%s", n_pages, use_page_split)
        except Exception:
            pass

    client = genai.Client(api_key=GEMINI_API_KEY)

    if not use_page_split:
        # Single-call path for short PDFs
        try:
            response = await client.aio.models.generate_content(
                model=GEMINI_OCR_MODEL,
                contents=[
                    genai_types.Part.from_bytes(data=pdf_bytes, mime_type="application/pdf"),
                    genai_types.Part.from_text(_OCR_PROMPT),
                ],
                config=genai_types.GenerateContentConfig(max_output_tokens=32768, thinking_config=genai_types.ThinkingConfig(thinking_budget=0)),
            )
            result = response.text or ""
            if on_progress:
                on_progress(1, 1, len(result))
            return result
        except Exception as e:
            logger.warning("[OCR-ASYNC] single-call failed: %s — trying page split", e)
            use_page_split = True

    # Page-by-page path: count pages first (no blob allocation), then iterate via generator.
    try:
        import io as _io2
        from pypdf import PdfReader as _PR2
        n_pages = len(_PR2(_io2.BytesIO(pdf_bytes)).pages)
    except Exception:
        n_pages = 0

    if n_pages == 0:
        # Fallback: send capped single call
        data = pdf_bytes[:_GEMINI_PDF_SIZE_LIMIT]
        try:
            response = await client.aio.models.generate_content(
                model=GEMINI_OCR_MODEL,
                contents=[
                    genai_types.Part.from_bytes(data=data, mime_type="application/pdf"),
                    genai_types.Part.from_text(_OCR_PROMPT),
                ],
                config=genai_types.GenerateContentConfig(max_output_tokens=32768, thinking_config=genai_types.ThinkingConfig(thinking_budget=0)),
            )
            return response.text or ""
        except Exception as e:
            logger.warning("[OCR-ASYNC] fallback call failed: %s", e)
            return ""

    pages_to_process = min(n_pages, max_pages)
    logger.info("[OCR-ASYNC] page-by-page: total=%d processing=%d", n_pages, pages_to_process)

    page_texts: list[str] = []
    chars_total = 0
    for idx, blob in enumerate(_iter_pdf_pages(pdf_bytes)):
        if idx >= max_pages:
            break
        try:
            response = await client.aio.models.generate_content(
                model=GEMINI_OCR_MODEL,
                contents=[
                    genai_types.Part.from_bytes(data=blob, mime_type="application/pdf"),
                    genai_types.Part.from_text(_OCR_PROMPT),
                ],
                config=genai_types.GenerateContentConfig(max_output_tokens=8192, thinking_config=genai_types.ThinkingConfig(thinking_budget=0)),
            )
            page_text = response.text or ""
            page_texts.append(page_text)
            chars_total += len(page_text)
        except Exception as e:
            logger.warning("[OCR-ASYNC] page %d failed: %s", idx + 1, e)
            page_texts.append("")

        if on_progress:
            on_progress(idx + 1, pages_to_process, chars_total)

    result = "\n\n".join(t for t in page_texts if t.strip())
    elapsed = (time.perf_counter() - t0) * 1000
    logger.info("[OCR-ASYNC] done: pages=%d elapsed=%.0fms chars=%d", len(capped), elapsed, len(result))
    if n_pages > max_pages:
        logger.warning("[OCR-ASYNC] PDF has %d pages, only first %d processed", n_pages, max_pages)
    return result


# ── xlsx / xls ────────────────────────────────────────────────────────────────

def _extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Convert Excel workbook sheets to Markdown tables for ETL processing."""
    try:
        import openpyxl
        from io import BytesIO

        wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            _XLSX_MAX_ROWS = 500  # per-sheet data-row cap
            if len(rows) > _XLSX_MAX_ROWS + 1:
                logger.warning(
                    "[OCR] xlsx: sheet=%r has %d data rows — only first %d extracted "
                    "(%d rows truncated)",
                    sheet_name, len(rows) - 1, _XLSX_MAX_ROWS, len(rows) - _XLSX_MAX_ROWS - 1,
                )
            parts.append(f"\n## Sheet: {sheet_name}\n")
            headers = [str(c) if c is not None else "" for c in rows[0]]
            n_cols = len(headers)
            parts.append("| " + " | ".join(headers) + " |")
            parts.append("| " + " | ".join(["---"] * n_cols) + " |")
            for row in rows[1: _XLSX_MAX_ROWS + 1]:
                raw = [str(c) if c is not None else "" for c in row]
                # Pad/truncate to match header width so the markdown table stays valid.
                padded = (raw + [""] * n_cols)[:n_cols]
                parts.append("| " + " | ".join(padded) + " |")
        wb.close()
        return "\n".join(parts)
    except Exception as e:
        logger.warning("[OCR] xlsx extraction failed: %s", e)
        return ""


# ── Entry point ───────────────────────────────────────────────────────────────

def extract_text_from_file(file_bytes: bytes, filename: str) -> tuple[str, str]:
    """
    Detect file format and extract text accordingly.

    Returns (extracted_text, detected_format).
    For scanned PDFs where text extraction fails, attempts Gemini Vision OCR.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    file_kb = len(file_bytes) // 1024
    logger.info("[OCR] extract_text_from_file: file=%r ext=%s size=%dKB", filename, ext, file_kb)
    t_start = time.perf_counter()

    if ext == "pdf":
        text = extract_text_from_pdf(file_bytes)
        if not text.strip():
            logger.info(
                "[OCR] PDF text extraction returned nothing (empty or low-quality) "
                "— escalating to Gemini Vision OCR: %r", filename,
            )
            text = extract_text_from_scanned_pdf_vision(file_bytes)
            if not text.strip():
                logger.warning("[OCR] Vision OCR also returned empty for %r", filename)
        fmt = "pdf"
    elif ext in ("docx",):
        text = extract_text_from_docx(file_bytes)
        fmt = "docx"
    elif ext in ("doc",):
        text = extract_text_from_docx(file_bytes)
        fmt = "doc"
    elif ext in ("pptx",):
        text = extract_text_from_pptx(file_bytes)
        fmt = "pptx"
    elif ext in ("ppt",):
        text = extract_text_from_pptx(file_bytes)
        fmt = "ppt"
    elif ext in ("txt", "csv", "md"):
        text = file_bytes.decode("utf-8", errors="replace")
        fmt = ext
    elif ext in ("jpg", "jpeg"):
        text = extract_text_from_image(file_bytes, "image/jpeg")
        fmt = "jpg"
    elif ext in ("png",):
        text = extract_text_from_image(file_bytes, "image/png")
        fmt = "png"
    elif ext in ("gif",):
        text = extract_text_from_image(file_bytes, "image/gif")
        fmt = "gif"
    elif ext in ("webp",):
        text = extract_text_from_image(file_bytes, "image/webp")
        fmt = "webp"
    elif ext in ("bmp",):
        text = extract_text_from_image(file_bytes, "image/bmp")
        fmt = "bmp"
    elif ext in ("tiff", "tif"):
        text = extract_text_from_image(file_bytes, "image/tiff")
        fmt = "tiff"
    elif ext == "xlsx":
        text = _extract_text_from_xlsx(file_bytes)
        fmt = "xlsx"
    elif ext == "xls":
        # openpyxl only supports .xlsx (Office Open XML); legacy .xls is unsupported.
        # Return a diagnostic string so ETL knows the file was received but unreadable.
        logger.warning("[OCR] .xls file uploaded (%r) — openpyxl only reads .xlsx; extraction skipped", filename)
        text = "[XLS_UNSUPPORTED: re-save as .xlsx to enable extraction]"
        fmt = "xls"
    else:
        logger.info("[OCR] unknown extension %r — trying pdf then docx", ext)
        text = extract_text_from_pdf(file_bytes)
        if not text:
            text = extract_text_from_docx(file_bytes)
        fmt = ext or "unknown"

    total_ms = (time.perf_counter() - t_start) * 1000
    stats = _quality_stats(text)
    logger.info(
        "[OCR] extract_text_from_file: DONE file=%r fmt=%s total_elapsed=%.0fms "
        "final_chars=%d meaningful=%d ratio=%.1f%%",
        filename, fmt, total_ms,
        stats["chars"], stats["meaningful"], stats["ratio_pct"],
    )
    if not text.strip():
        logger.warning(
            "[OCR] extract_text_from_file: FINAL TEXT IS EMPTY for %r — "
            "ETL will receive no text and return no data", filename,
        )
    return text, fmt
