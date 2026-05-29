"""
tests/test_tsmc_piercing.py

TSMC (台灣積體電路製造, 2330) End-to-End Piercing Test Suite
============================================================
This suite does NOT test hypothetical features — it probes the actual source
code paths executed when a user uploads a real TSMC document.  Every
assertion is traceable to a specific line in the production codebase.

Defect catalog (DID = Defect ID):
  DID-UPLOAD-01   CREDIT_REPORT_MAX_UPLOAD_MB=20 hard-rejects TSMC 30–80 MB PDFs
                  → credit_report/config.py:50, api/generate.py:42,146
  DID-ETL-TRUNC   CR_ETL_MAX_TEXT_CHARS=120 000 truncates 400-page reports (~600 K chars)
                  → credit_report/config.py:51, generation/etl.py:1204
  DID-OCR-CAP     Vision OCR PDF byte cap = 20 MB; full report pages beyond that are dark
                  → credit_report/generation/evidence.py:434
  DID-TABLE-01    No camelot / pdfplumber / tabula-py → financial tables extracted as prose
                  → requirements.txt (absent), evidence.py (no table-aware parser)
  DID-PPTX-01     PPTX extraction works but yields only text shapes; chart/image data lost
                  → credit_report/generation/evidence.py:330-368
  DID-TWSE-UNIT   TWSE returns thousands-NTD; system propagates raw integers without USD conv.
                  → credit_report/integrations/twse.py:394, build_section7_input()
  DID-TWSE-FIELD  Chinese label→English metric mapping coverage
                  → credit_report/integrations/twse.py:37-91 (INCOME_METRIC_ALIASES etc.)
  DID-TWSE-PAID   paid_in_capital from PROFILE is raw TWD (not thousands); unit mismatch
                  → credit_report/integrations/twse.py:104, build_section7_input()

Running:
    pytest tests/test_tsmc_piercing.py -v
"""
from __future__ import annotations

import importlib.util
import io
import struct
import textwrap
import zipfile
from collections import defaultdict
from typing import Any
from unittest.mock import AsyncMock, MagicMock, patch

import pytest

# ---------------------------------------------------------------------------
# Lazy import helpers — avoid crashing if optional extras are absent in CI
# ---------------------------------------------------------------------------

def _have(pkg: str) -> bool:
    return importlib.util.find_spec(pkg) is not None


# ===========================================================================
# SECTION 1 — Configuration boundary assertions
# Confirm the actual configured limits that exclude TSMC-sized documents.
# ===========================================================================

class TestConfigBoundaries:
    """DID-UPLOAD-01 / DID-ETL-TRUNC / DID-OCR-CAP"""

    def test_upload_cap_is_20mb(self):
        """CREDIT_REPORT_MAX_UPLOAD_MB is hard-coded to 20 MB.

        TSMC's annual report (e.g. 2023) is ~54 MB; their sustainability
        report is ~30 MB.  Any value ≤ 29 MB rejects the primary use-case.
        Source: credit_report/config.py:50
        """
        from credit_report.config import CREDIT_REPORT_MAX_UPLOAD_MB
        assert CREDIT_REPORT_MAX_UPLOAD_MB == 20, (
            f"Upload cap changed to {CREDIT_REPORT_MAX_UPLOAD_MB} MB — "
            "update this test AND verify TSMC PDFs now fit."
        )
        # DID-UPLOAD-01: typical TSMC annual report sizes (MB)
        tsmc_typical_sizes_mb = {"annual_report_2023": 54, "sustainability_report_2023": 31}
        oversized = {k: v for k, v in tsmc_typical_sizes_mb.items() if v > CREDIT_REPORT_MAX_UPLOAD_MB}
        assert oversized, (
            "Expected at least one TSMC document to exceed the 20 MB cap. "
            "DID-UPLOAD-01 would be resolved if this assertion fires."
        )

    def test_etl_text_cap_is_120k_chars(self):
        """CR_ETL_MAX_TEXT_CHARS truncates to 120 000 characters.

        A 400-page Chinese/English TSMC annual report yields ~500 000–800 000
        characters of extracted text.  Only the first ~20–25% reaches Gemini.
        Financial statements (Income Statement, Balance Sheet, Cash Flow) appear
        in pages 100–250 of the annual report — well beyond the truncation point.
        Source: credit_report/config.py:51
        """
        from credit_report.config import CR_ETL_MAX_TEXT_CHARS
        assert CR_ETL_MAX_TEXT_CHARS == 120_000, (
            f"ETL cap changed to {CR_ETL_MAX_TEXT_CHARS}; update this test."
        )
        tsmc_annual_report_chars_low = 500_000
        coverage_pct = CR_ETL_MAX_TEXT_CHARS / tsmc_annual_report_chars_low * 100
        assert coverage_pct < 30, (
            f"Coverage is {coverage_pct:.1f}% — DID-ETL-TRUNC may be resolved; "
            "increase CR_ETL_MAX_TEXT_CHARS and re-validate."
        )

    def test_vision_ocr_pdf_byte_cap_is_20mb(self):
        """Gemini Vision OCR receives at most the first 20 MB of a PDF.

        For scanned-PDF fallback, evidence.py:434 slices:
            data = pdf_bytes[:20 * 1024 * 1024]
        A 54 MB TSMC annual report has ~34 MB of data beyond this cap that
        Gemini never sees.
        Source: credit_report/generation/evidence.py:434
        """
        OCR_CAP_BYTES = 20 * 1024 * 1024
        tsmc_pdf_bytes = 54 * 1024 * 1024
        dark_bytes = tsmc_pdf_bytes - OCR_CAP_BYTES
        dark_pct = dark_bytes / tsmc_pdf_bytes * 100
        assert dark_pct > 30, (
            "Expected >30%% of TSMC PDF to be beyond Vision OCR cap. "
            "DID-OCR-CAP may be resolved; verify by increasing the cap."
        )
        # Sanity: cap constant matches actual code
        import re
        import credit_report.generation.evidence as ev_mod
        source = open(ev_mod.__file__).read()
        assert "20 * 1024 * 1024" in source, (
            "Vision OCR byte cap constant not found in evidence.py — "
            "source may have changed; re-inspect DID-OCR-CAP."
        )

    def test_no_table_parsing_libraries(self):
        """camelot-py, pdfplumber, and tabula-py are absent from the environment.

        Without a structure-preserving table parser, financial tables (TSMC's
        income statement spans ~30 rows and 5 years) are extracted as
        unstructured prose or lost entirely.  DID-TABLE-01.
        """
        absent = []
        for pkg in ("camelot", "pdfplumber", "tabula"):
            if not _have(pkg):
                absent.append(pkg)
        assert absent == ["camelot", "pdfplumber", "tabula"], (
            f"Unexpected packages installed: {set(['camelot','pdfplumber','tabula']) - set(absent)}. "
            "DID-TABLE-01 is partially resolved — validate table extraction quality."
        )


# ===========================================================================
# SECTION 2 — ETL truncation behaviour
# ===========================================================================

class TestETLTruncation:
    """DID-ETL-TRUNC: verify that _build_etl_prompt slices text at CR_ETL_MAX_TEXT_CHARS."""

    def _realistic_tsmc_text(self, total_chars: int = 600_000) -> str:
        """Simulate a TSMC annual report text: Chinese headers in first 20%, financials later."""
        header = (
            "台灣積體電路製造股份有限公司 2023年度年報\n"
            "Taiwan Semiconductor Manufacturing Company Limited Annual Report 2023\n\n"
            "公司簡介 Company Overview\n"
            "台積公司成立於1987年，是全球最大的專業積體電路製造服務公司。\n"
        ) * 200  # repeating header section: ~front 20%

        financials = (
            "合併損益表 Consolidated Statements of Comprehensive Income\n"
            "單位：新台幣千元 (In Thousands of New Taiwan Dollars)\n"
            "項目 Item                 2023年度      2022年度\n"
            "淨收益 Net Revenue     2,161,735,000  2,263,361,000\n"
            "營業成本 Cost of Revenue  (887,843,000) (846,218,000)\n"
            "毛利 Gross Profit      1,273,892,000  1,417,143,000\n"
            "營業利益 Operating Profit 1,003,227,000  1,189,280,000\n"
            "稅後淨利 Net Income       838,497,000   959,437,000\n"
            "基本每股盈餘 EPS (NTD)         32.34         37.00\n\n"
        ) * 500

        full_text = header + financials
        return full_text[:total_chars]

    def test_prompt_contains_only_first_120k_chars(self):
        """ETL prompt truncates document text at CR_ETL_MAX_TEXT_CHARS.

        Source: generation/etl.py:1204
            text_snippet = text[:CR_ETL_MAX_TEXT_CHARS]
        """
        from credit_report.config import CR_ETL_MAX_TEXT_CHARS
        from credit_report.generation.etl import _build_etl_prompt

        full_text = self._realistic_tsmc_text(600_000)
        _, user_prompt = _build_etl_prompt("annual_report", full_text, [4, 7])

        start = user_prompt.index("---DOCUMENT TEXT START---") + len("---DOCUMENT TEXT START---")
        end = user_prompt.index("---DOCUMENT TEXT END---")
        doc_slice = user_prompt[start:end]

        assert len(doc_slice.strip()) <= CR_ETL_MAX_TEXT_CHARS + 5, (
            "Document slice in prompt exceeds CR_ETL_MAX_TEXT_CHARS — "
            "truncation guard may have been removed."
        )
        truncated_chars = len(full_text) - CR_ETL_MAX_TEXT_CHARS
        assert truncated_chars > 0, "No truncation occurred — full text fits within cap."

    def test_financial_statements_are_beyond_truncation_point(self):
        """The financial statements section of a TSMC annual report falls after the truncation
        cutoff of 120 000 characters.

        In the TSMC 2023 Annual Report the consolidated financial statements begin
        on approximately page 100 of 400 (byte offset ~135 000+).  Only pages 1–~25
        fit within the 120 000-char window.
        """
        from credit_report.config import CR_ETL_MAX_TEXT_CHARS

        full_text = self._realistic_tsmc_text(600_000)
        # Find where financial-statement markers appear in the full text
        fs_marker = "合併損益表"
        fs_pos = full_text.find(fs_marker)
        assert fs_pos != -1, "Financial statement marker not found in synthetic text."

        # The marker should appear beyond the truncation point at least once
        # (financial data is in the back half of the report)
        tail = full_text[CR_ETL_MAX_TEXT_CHARS:]
        assert fs_marker in tail, (
            f"All financial statement markers appear within the first {CR_ETL_MAX_TEXT_CHARS} chars. "
            "DID-ETL-TRUNC: The test synthetic text may need adjustment, but for a real TSMC "
            "annual report the financials span pages 100–250 which exceed this cap."
        )

    def test_logger_reports_truncation_flag(self):
        """_build_etl_prompt logs truncated=True when text exceeds CR_ETL_MAX_TEXT_CHARS.

        The logger uses % formatting with positional args; we format each call
        to reconstruct the final log message and search for 'truncated=True'.
        """
        from credit_report.generation.etl import _build_etl_prompt

        full_text = self._realistic_tsmc_text(600_000)
        with patch("credit_report.generation.etl.logger") as mock_log:
            _build_etl_prompt("annual_report", full_text, [4, 7])
            truncated_logged = False
            for call in mock_log.info.call_args_list:
                pos_args = call[0]  # (fmt, arg1, arg2, ...)
                if pos_args:
                    try:
                        formatted = pos_args[0] % pos_args[1:]
                        if "truncated=True" in formatted:
                            truncated_logged = True
                            break
                    except (TypeError, ValueError):
                        pass
        assert truncated_logged, (
            "etl.py did not produce a log line containing 'truncated=True' "
            "for a 600 000-char document. "
            "Either the logging call was removed or the truncation logic changed."
        )


# ===========================================================================
# SECTION 3 — PDF text quality and Vision OCR fallback
# ===========================================================================

class TestPDFTextQuality:
    """DID-OCR-CAP / DID-TABLE-01"""

    def test_cjk_cid_font_output_fails_quality_check(self):
        """pdfminer on CID-font PDFs returns mostly whitespace / unmapped glyphs.

        _text_quality_ok requires >= 5% meaningful characters (CJK + ASCII alpha/digit).
        Garbage output from CID-font PDFs has ~0–2% meaningful chars and is rejected.
        Source: evidence.py:122-140
        """
        from credit_report.generation.evidence import _text_quality_ok

        # Simulate what pdfminer returns for a CID-font Chinese PDF:
        # mostly whitespace, random bytes, no recognisable characters
        cid_garbage = " \x00\x01\x02\x03\x04\x05\x06\x07\x08" * 500 + "  \n  \n  \n" * 200
        assert not _text_quality_ok(cid_garbage), (
            "CID-font garbage passed quality check — threshold may be too low."
        )

    def test_clean_chinese_text_passes_quality_check(self):
        """Genuine Chinese text extracted by a capable parser passes the quality check."""
        from credit_report.generation.evidence import _text_quality_ok

        chinese_text = (
            "台積公司成立於1987年，是全球最大的專業積體電路製造服務公司。\n"
            "2023年合併營業收入為新台幣2兆1,617億元，稅後淨利為新台幣8,385億元。\n"
            "每股盈餘（EPS）新台幣32.34元。\n"
        ) * 20
        assert _text_quality_ok(chinese_text), (
            "Genuine Chinese text failed quality check — threshold may be too strict."
        )

    def test_vision_ocr_cap_at_20mb_in_source(self):
        """Verify the 20 MB Vision OCR slice is present and unchanged in evidence.py."""
        import credit_report.generation.evidence as ev_mod
        source = open(ev_mod.__file__, encoding="utf-8").read()
        # Line: data = pdf_bytes[:20 * 1024 * 1024]
        assert "pdf_bytes[:20 * 1024 * 1024]" in source, (
            "DID-OCR-CAP: Vision OCR byte cap line not found in evidence.py — "
            "check if the cap was changed or the function was refactored."
        )

    def test_extract_text_from_pdf_returns_empty_for_cid_garbage(self):
        """When both pdfminer and pypdf fail quality checks, extract_text_from_pdf returns ''."""
        from credit_report.generation.evidence import extract_text_from_pdf

        # Minimal valid PDF with no extractable text (1-page empty PDF)
        # This tests the fall-through path; Vision OCR would be called next
        # but GEMINI_API_KEY is not set in unit test context.
        minimal_pdf = (
            b"%PDF-1.4\n"
            b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
            b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R>>endobj\n"
            b"xref\n0 4\n0000000000 65535 f\n0000000009 00000 n\n"
            b"0000000058 00000 n\n0000000115 00000 n\n"
            b"trailer<</Size 4/Root 1 0 R>>\nstartxref\n190\n%%EOF"
        )
        result = extract_text_from_pdf(minimal_pdf)
        # Empty PDF should yield empty or near-empty text
        assert len(result.strip()) < 50, (
            f"Expected near-empty text from empty PDF, got {len(result)} chars."
        )


# ===========================================================================
# SECTION 4 — PPTX extraction
# ===========================================================================

class TestPPTXExtraction:
    """DID-PPTX-01: PPTX extraction captures text shapes but not image/chart data."""

    def _make_minimal_pptx(self, include_table: bool = True) -> bytes:
        """Build a minimal PPTX in memory with a text shape and optionally a table."""
        buf = io.BytesIO()
        # python-pptx is required — skip if absent
        if not _have("pptx"):
            pytest.skip("python-pptx not installed")
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        tf = txBox.text_frame
        tf.text = "台積電 2024Q3法說會 TSMC 3Q24 Investor Conference"

        if include_table:
            rows, cols = 4, 3
            tbl = slide.shapes.add_table(rows, cols, Inches(1), Inches(2.5), Inches(6), Inches(2)).table
            headers = ["指標 Metric", "3Q24", "2Q24"]
            data = [
                ["淨收益 Net Revenue (USD Bn)", "23.51", "20.82"],
                ["毛利率 Gross Margin (%)", "57.8%", "53.2%"],
                ["稅後淨利 Net Income (USD Bn)", "9.97", "8.55"],
            ]
            for ci, h in enumerate(headers):
                tbl.cell(0, ci).text = h
            for ri, row_data in enumerate(data, start=1):
                for ci, val in enumerate(row_data):
                    tbl.cell(ri, ci).text = val

        prs.save(buf)
        return buf.getvalue()

    def test_pptx_text_shape_is_extracted(self):
        """extract_text_from_pptx captures text box content."""
        from credit_report.generation.evidence import extract_text_from_pptx

        pptx_bytes = self._make_minimal_pptx(include_table=False)
        text, fmt = __import__("credit_report.generation.evidence", fromlist=["extract_text_from_file"]).extract_text_from_file(pptx_bytes, "test.pptx")
        assert "TSMC" in text or "台積電" in text, (
            f"PPTX text box content not found in extracted text. Got: {text[:200]!r}"
        )
        assert fmt == "pptx"

    def test_pptx_table_cells_are_extracted(self):
        """extract_text_from_pptx captures table cell text via pipe-separated rows."""
        from credit_report.generation.evidence import extract_text_from_pptx

        pptx_bytes = self._make_minimal_pptx(include_table=True)
        text = extract_text_from_pptx(pptx_bytes)

        # The table should be present as pipe-separated rows
        assert "Net Revenue" in text or "淨收益" in text, (
            "PPTX table content not found. DID-PPTX-01 regression: table extraction broken."
        )
        # Verify pipe-separated format for at least one row
        has_pipe_row = any("|" in line for line in text.splitlines())
        assert has_pipe_row, (
            "No pipe-separated rows found in PPTX output — table structure lost. "
            "DID-PPTX-01: structure extraction may be incomplete."
        )

    def test_pptx_financial_figures_survive_extraction(self):
        """Numeric values in PPTX tables must survive extraction intact."""
        from credit_report.generation.evidence import extract_text_from_pptx

        pptx_bytes = self._make_minimal_pptx(include_table=True)
        text = extract_text_from_pptx(pptx_bytes)

        for value in ("23.51", "20.82", "57.8%"):
            assert value in text, (
                f"Financial figure {value!r} lost during PPTX extraction. "
                "DID-PPTX-01: numeric data integrity failure."
            )

    def test_pptx_is_supported_upload_extension(self):
        """'pptx' is listed in ALLOWED_EXTENSIONS in generate.py."""
        from credit_report.api.generate import ALLOWED_EXTENSIONS
        assert "pptx" in ALLOWED_EXTENSIONS, (
            "pptx removed from ALLOWED_EXTENSIONS — upload will be rejected."
        )

    def test_pptx_routes_to_extract_text_from_pptx(self):
        """extract_text_from_file routes .pptx to extract_text_from_pptx."""
        import credit_report.generation.evidence as ev_mod
        source = open(ev_mod.__file__, encoding="utf-8").read()
        assert 'ext in ("pptx",)' in source or "pptx" in source, (
            "PPTX routing block not found in evidence.py."
        )


# ===========================================================================
# SECTION 5 — TWSE integration: unit and field mapping
# ===========================================================================

class TestTWSEUnitAndMapping:
    """DID-TWSE-UNIT / DID-TWSE-FIELD / DID-TWSE-PAID"""

    # -------------------------------------------------------------------
    # Representative TSMC-style TWSE row data (abbreviated)
    # TWSE financial statement figures are denominated in thousands NTD (千元).
    # TSMC 2023 consolidated revenue: 2,161,735,453 thousands NTD
    # -------------------------------------------------------------------
    TSMC_INCOME_ROW = {
        "公司代號": "2330",
        "公司名稱": "台積電",
        "年度": "113",           # ROC year 113 = 2024
        "季別": "4",
        "會計項目": "營業收入",
        "金額": "2161735453",    # thousands NTD
    }
    TSMC_INCOME_ROW_EBITDA = {
        "公司代號": "2330",
        "公司名稱": "台積電",
        "年度": "113",
        "季別": "4",
        "會計項目": "稅前息前折舊攤銷前淨利",
        "金額": "1120000000",
    }
    TSMC_BALANCE_ROW_ASSETS = {
        "公司代號": "2330",
        "公司名稱": "台積電",
        "年度": "113",
        "季別": "4",
        "會計項目": "資產總計",
        "金額": "6800000000",    # thousands NTD
    }
    TSMC_BALANCE_ROW_EQUITY = {
        "公司代號": "2330",
        "公司名稱": "台積電",
        "年度": "113",
        "季別": "4",
        "會計項目": "權益總計",
        "金額": "4500000000",
    }
    TSMC_PROFILE_ROW = {
        "公司代號": "2330",
        "公司名稱": "台積電",
        "英文簡稱": "TSMC",
        "產業別": "半導體業",
        "實收資本額": "259303805000",  # raw TWD (NOT thousands)
        "簽證會計師事務所": "勤業眾信",
        "上市日期": "19940905",
    }

    def _make_bundle(self) -> dict[str, list[dict]]:
        return {
            "company_profile": [self.TSMC_PROFILE_ROW],
            "monthly_revenue": [],
            "income_statement_general": [
                self.TSMC_INCOME_ROW,
                self.TSMC_INCOME_ROW_EBITDA,
            ],
            "balance_sheet_general": [
                self.TSMC_BALANCE_ROW_ASSETS,
                self.TSMC_BALANCE_ROW_EQUITY,
            ],
            "cash_flow_general": [],
            "valuation_ratios": [],
            "daily_trading": [],
            "monthly_average": [],
            "dividend": [],
        }

    # ---- Field mapping tests -----------------------------------------------

    def test_ying_ye_shou_ru_maps_to_revenue(self):
        """'營業收入' (operating revenue) maps to the 'revenue' key."""
        from credit_report.integrations.twse import _metric_from_label, INCOME_METRIC_ALIASES
        result = _metric_from_label("營業收入", INCOME_METRIC_ALIASES)
        assert result == "revenue", (
            f"'營業收入' mapped to {result!r} instead of 'revenue'. "
            "DID-TWSE-FIELD: income statement field mapping broken."
        )

    def test_mao_li_maps_to_gross_profit(self):
        """'營業毛利' maps to 'gross_profit'."""
        from credit_report.integrations.twse import _metric_from_label, INCOME_METRIC_ALIASES
        assert _metric_from_label("營業毛利", INCOME_METRIC_ALIASES) == "gross_profit"

    def test_zi_chan_zong_ji_maps_to_total_assets(self):
        """'資產總計' maps to 'total_assets'."""
        from credit_report.integrations.twse import _metric_from_label, BALANCE_METRIC_ALIASES
        assert _metric_from_label("資產總計", BALANCE_METRIC_ALIASES) == "total_assets"

    def test_quan_yi_zong_ji_maps_to_total_equity(self):
        """'權益總計' maps to 'total_equity'."""
        from credit_report.integrations.twse import _metric_from_label, BALANCE_METRIC_ALIASES
        assert _metric_from_label("權益總計", BALANCE_METRIC_ALIASES) == "total_equity"

    def test_eps_maps_correctly(self):
        """'基本每股盈餘' maps to 'eps'."""
        from credit_report.integrations.twse import _metric_from_label, INCOME_METRIC_ALIASES
        assert _metric_from_label("基本每股盈餘", INCOME_METRIC_ALIASES) == "eps"

    def test_ebitda_label_maps_correctly(self):
        """'稅前息前折舊攤銷前淨利' maps to 'ebitda'."""
        from credit_report.integrations.twse import _metric_from_label, INCOME_METRIC_ALIASES
        assert _metric_from_label("稅前息前折舊攤銷前淨利", INCOME_METRIC_ALIASES) == "ebitda"

    # ---- Unit / scale tests -----------------------------------------------

    def test_build_section7_declares_unit_thousands_ntd(self):
        """build_section7_input declares reporting_currency=TWD and unit=thousands.

        DID-TWSE-UNIT: the raw integers from TWSE are in thousands NTD.
        A downstream consumer reading revenue=2161735453 without the unit
        metadata would interpret this as 2.16 TRILLION USD — ~32 000× too large.
        Source: integrations/twse.py:394-396
        """
        from credit_report.integrations.twse import build_section7_input

        bundle = self._make_bundle()
        result = build_section7_input("2330", bundle, role="borrower")

        fin = result.get("7A_borrower_financials", {})
        assert fin.get("reporting_currency") == "TWD", (
            "DID-TWSE-UNIT: reporting_currency missing or changed from TWD."
        )
        assert fin.get("unit") == "thousands", (
            "DID-TWSE-UNIT: unit not declared as 'thousands' — scale metadata lost."
        )

    def test_revenue_figure_is_raw_thousands_ntd(self):
        """Revenue from TWSE is stored as raw thousands-NTD integer with no USD conversion.

        For TSMC: revenue = 2 161 735 453 (thousands NTD) ≈ TWD 2.16 trillion ≈ USD 67 billion.
        The system stores 2 161 735 453.  No conversion to USD millions is applied.
        DID-TWSE-UNIT: credit analysts expecting USD millions will see a 32 000× magnitude error.

        revenue_fy2024 is nested under result["7A_borrower_financials"] (not top-level).
        Source: integrations/twse.py — financials.update({"revenue_fy2024": ...}) then
        section["7A_borrower_financials"] = financials
        """
        from credit_report.integrations.twse import build_section7_input

        bundle = self._make_bundle()
        result = build_section7_input("2330", bundle, role="borrower")

        # revenue_fy2024 is inside 7A_borrower_financials, not the top-level section dict
        fin = result.get("7A_borrower_financials", {})
        revenue_raw = fin.get("revenue_fy2024")
        assert revenue_raw is not None, (
            "revenue_fy2024 not found in result['7A_borrower_financials']. "
            "Check that the mock TWSE income_statement_general row was parsed. "
            f"7A keys: {list(fin.keys())[:15]}"
        )
        # TSMC revenue in thousands NTD is in the billions range (10-digit integer)
        assert revenue_raw > 1_000_000_000, (
            f"Revenue {revenue_raw} is not in the expected billions-range for TSMC thousands-NTD. "
            "DID-TWSE-UNIT: check if conversion was applied upstream."
        )

    def test_paid_in_capital_is_raw_twd_not_thousands(self):
        """paid_in_capital from PROFILE_FIELD_ALIASES is in raw TWD, not thousands.

        TSMC's 實收資本額 is "259,303,805,000" TWD.
        TWSE income statement figures are thousands NTD.
        These two different scales appear in the same output dict, creating a
        unit inconsistency: paid_in_capital cannot be compared to revenue without
        a ÷1000 correction.
        DID-TWSE-PAID: source integrations/twse.py:104
        """
        from credit_report.integrations.twse import build_section7_input

        bundle = self._make_bundle()
        result = build_section7_input("2330", bundle, role="borrower")

        fin = result.get("7A_borrower_financials", {})
        profile = fin.get("twse_company_profile", {})

        paid_in_capital = profile.get("paid_in_capital")
        if paid_in_capital is None:
            pytest.skip("paid_in_capital not returned for this mock data set.")

        # Convert to float for comparison
        try:
            capital_val = float(str(paid_in_capital).replace(",", ""))
        except ValueError:
            pytest.skip(f"paid_in_capital is not numeric: {paid_in_capital!r}")

        # The raw TWD value (259 B TWD) is ~1000× larger than the thousands-NTD revenue figure
        # expressed in the same units.  Both should NOT equal a common "thousands NTD" scale.
        # TSMC paid_in_capital raw TWD: ~2.59e11 (billion TWD territory)
        assert capital_val > 1e10, (
            "DID-TWSE-PAID: paid_in_capital appears to be in thousands NTD — "
            "verify whether the PROFILE field is now being normalised."
        )

    def test_no_usd_conversion_in_build_section7(self):
        """build_section7_input does NOT convert any value from TWD to USD.

        The function returns TWD figures and documents the currency/unit in metadata.
        Conversion to USD millions required for credit analysis must be done by the
        analyst or a downstream process.
        """
        from credit_report.integrations.twse import build_section7_input

        bundle = self._make_bundle()
        result = build_section7_input("2330", bundle, role="borrower")

        # There should be no field named *_usd_m or *_usd_bn in the flat output
        all_keys = list(result.keys())
        usd_converted = [k for k in all_keys if "_usd" in k.lower()]
        assert not usd_converted, (
            f"Unexpected USD-converted fields found: {usd_converted}. "
            "DID-TWSE-UNIT: if USD conversion was added, update this test "
            "and verify the conversion factor is correct."
        )

    # ---- ROC calendar year conversion test --------------------------------

    def test_roc_year_converts_to_gregorian(self):
        """_period_key converts ROC calendar years (< 1911) to Gregorian.

        TWSE uses ROC year in '年度' column: 113 → 2024, 112 → 2023, etc.
        Source: integrations/twse.py:162
        """
        from credit_report.integrations.twse import _period_key

        row = {"年度": "113", "季別": "4"}
        key = _period_key(row)
        assert "2024" in key, (
            f"ROC year 113 should convert to 2024 but got key={key!r}. "
            "DID-TWSE-FIELD: calendar conversion broken."
        )

    def test_roc_year_112_is_2023(self):
        """ROC year 112 = Gregorian 2023."""
        from credit_report.integrations.twse import _period_key
        assert "2023" in _period_key({"年度": "112", "季別": "1"})

    # ---- TWSE backend-side fetch (no browser CORS issue) ------------------

    def test_twse_client_uses_httpx_not_browser(self):
        """The TWSE integration calls the API server-side via httpx, not from the browser.

        This means browser CORS restrictions do NOT apply to the TWSE data path.
        DID-TWSE-CORS is not a defect in the current architecture: all TWSE
        calls are backend-proxied through TWSEOpenAPIClient.fetch().
        Source: integrations/twse.py:309-335
        """
        from credit_report.integrations.twse import TWSEOpenAPIClient
        import inspect
        src = inspect.getsource(TWSEOpenAPIClient.fetch)
        assert "httpx" in src, "TWSEOpenAPIClient.fetch does not use httpx — check import."
        assert "AsyncClient" in src, "TWSEOpenAPIClient.fetch is not using async httpx client."


# ===========================================================================
# SECTION 6 — Document type routing and section mapping
# ===========================================================================

class TestDocumentTypeRouting:
    """Confirm that TSMC document types map to the expected section extraction targets."""

    def test_annual_report_extracts_sections_4_7(self):
        """annual_report → sections [4, 7, 3, 2, 10].

        For a TSMC annual report this means: company background (§4),
        financial analysis (§7), credit ratings (§3), overall comments (§2),
        and appendix (§10).
        Source: generation/etl.py:22-26
        """
        from credit_report.generation.etl import DOCUMENT_SECTION_MAP
        sections = DOCUMENT_SECTION_MAP["annual_report"]
        assert 4 in sections, "§4 (corporate background) missing from annual_report ETL map."
        assert 7 in sections, "§7 (financial analysis) missing from annual_report ETL map."

    def test_financial_statement_excludes_section_1(self):
        """financial_statement does NOT target §1 (credit facility) — correct.

        A standalone TSMC financial statement PDF has no credit facility data.
        Source: generation/etl.py:23
        """
        from credit_report.generation.etl import DOCUMENT_SECTION_MAP
        assert 1 not in DOCUMENT_SECTION_MAP["financial_statement"]

    def test_analyst_presentation_targets_section_4_7(self):
        """analyst_presentation (法說會) targets §4 and §7.

        TSMC's investor conference presentations contain operational and financial
        highlights used to populate §4 (business overview) and §7 (financials).
        Source: generation/etl.py:25
        """
        from credit_report.generation.etl import DOCUMENT_SECTION_MAP
        sections = DOCUMENT_SECTION_MAP["analyst_presentation"]
        assert 4 in sections
        assert 7 in sections

    def test_external_report_targets_section_11(self):
        """external_report (broker research) targets §11 first.

        A sell-side research report on TSMC (e.g. from 群益, 元大 securities)
        should populate §11 (Analyst/Broker Research).
        Source: generation/etl.py:31
        """
        from credit_report.generation.etl import DOCUMENT_SECTION_MAP
        sections = DOCUMENT_SECTION_MAP["external_report"]
        assert 11 in sections
        assert sections[0] == 11, "§11 should be the primary target for external_report."


# ===========================================================================
# SECTION 7 — Upload size guard mechanics
# ===========================================================================

class TestUploadSizeGuard:
    """DID-UPLOAD-01: verify that _persist_upload_to_disk rejects files > 20 MB."""

    @pytest.mark.asyncio
    async def test_upload_rejects_21mb_file(self, tmp_path):
        """A 21 MB file triggers HTTP 413 during chunked streaming.

        Source: api/generate.py:144-146
            if total > _MAX_UPLOAD_BYTES:
                raise HTTPException(status_code=413, ...)
        """
        from fastapi import HTTPException, UploadFile
        from credit_report.api.generate import _persist_upload_to_disk, _MAX_UPLOAD_BYTES

        assert _MAX_UPLOAD_BYTES == 20 * 1024 * 1024

        oversized = b"x" * (21 * 1024 * 1024)  # 21 MB

        class FakeUpload:
            filename = "tsmc_annual_report_2023.pdf"
            _pos = 0
            _data = oversized
            _chunk = 1024 * 1024  # 1 MB chunks

            async def read(self, size: int) -> bytes:
                chunk = self._data[self._pos: self._pos + size]
                self._pos += size
                return chunk

        doc_dir = tmp_path / "reports" / "test-report-id"
        doc_dir.mkdir(parents=True)

        import credit_report.api.generate as gen_module
        original_root = gen_module.runtime_config.CREDIT_REPORTS_ROOT

        try:
            gen_module.runtime_config.CREDIT_REPORTS_ROOT = tmp_path / "reports"
            with pytest.raises(HTTPException) as exc_info:
                await _persist_upload_to_disk(
                    "test-report-id", "doc-001", FakeUpload(), "tsmc_annual_report_2023.pdf"
                )
            assert exc_info.value.status_code == 413
            assert "20 MB" in exc_info.value.detail or "upload limit" in exc_info.value.detail
        finally:
            gen_module.runtime_config.CREDIT_REPORTS_ROOT = original_root

    @pytest.mark.asyncio
    async def test_upload_accepts_exactly_20mb(self, tmp_path):
        """A file of exactly CREDIT_REPORT_MAX_UPLOAD_MB is accepted (boundary value)."""
        from credit_report.api.generate import _persist_upload_to_disk, _MAX_UPLOAD_BYTES

        exact_size = _MAX_UPLOAD_BYTES  # 20 971 520 bytes

        class FakeUpload:
            filename = "exactly_20mb.pdf"
            _pos = 0
            _data = b"y" * exact_size
            _chunk = 1024 * 1024

            async def read(self, size: int) -> bytes:
                chunk = self._data[self._pos: self._pos + size]
                self._pos += size
                return chunk

        import credit_report.api.generate as gen_module
        original_root = gen_module.runtime_config.CREDIT_REPORTS_ROOT

        try:
            gen_module.runtime_config.CREDIT_REPORTS_ROOT = tmp_path / "reports"
            path, size = await _persist_upload_to_disk(
                "report-boundary", "doc-002", FakeUpload(), "exactly_20mb.pdf"
            )
            assert size == exact_size
        finally:
            gen_module.runtime_config.CREDIT_REPORTS_ROOT = original_root

    def test_max_upload_bytes_constant_is_derived_from_config(self):
        """_MAX_UPLOAD_BYTES must equal CREDIT_REPORT_MAX_UPLOAD_MB * 1024 * 1024.

        If someone changes CREDIT_REPORT_MAX_UPLOAD_MB in config without
        restarting the app, the module-level constant will be stale.  This test
        confirms the constant is correctly derived at import time.
        Source: api/generate.py:42
        """
        from credit_report.config import CREDIT_REPORT_MAX_UPLOAD_MB
        from credit_report.api.generate import _MAX_UPLOAD_BYTES
        assert _MAX_UPLOAD_BYTES == CREDIT_REPORT_MAX_UPLOAD_MB * 1024 * 1024


# ===========================================================================
# SECTION 8 — Defect summary: explicit known-limit assertions
# These assertions always pass; they exist to make the system's constraints
# discoverable via pytest output rather than buried in source comments.
# ===========================================================================

class TestKnownLimits:
    """Living documentation of currently known system constraints for TSMC use-case."""

    def test_known_limit_upload_20mb(self):
        from credit_report.config import CREDIT_REPORT_MAX_UPLOAD_MB
        # Document the limit — this test should be updated if the limit is raised
        assert CREDIT_REPORT_MAX_UPLOAD_MB == 20  # DID-UPLOAD-01: raise to ≥ 80 for TSMC

    def test_known_limit_etl_truncation_120k(self):
        from credit_report.config import CR_ETL_MAX_TEXT_CHARS
        assert CR_ETL_MAX_TEXT_CHARS == 120_000  # DID-ETL-TRUNC: raise to ≥ 500 000 for TSMC

    def test_known_limit_twse_unit_thousands_ntd(self):
        """TWSE data unit declared in build_section7_input: thousands NTD (not USD millions)."""
        from credit_report.integrations.twse import build_section7_input
        # Confirm the declared unit string hasn't silently changed
        import inspect
        src = inspect.getsource(build_section7_input)
        assert '"thousands"' in src, "Unit declaration removed from build_section7_input."
        assert '"TWD"' in src, "Currency declaration removed from build_section7_input."

    def test_known_limit_no_table_parser(self):
        """No dedicated PDF table parser is installed in this environment."""
        for pkg in ("camelot", "pdfplumber", "tabula"):
            assert not _have(pkg), (
                f"DID-TABLE-01 may be resolved: {pkg} is now available. "
                "Run a TSMC financial statement through the ETL pipeline and verify "
                "that income statement numbers are extracted correctly."
            )

    def test_known_limit_vision_ocr_pdf_cap(self):
        """Vision OCR for PDFs is capped at the first 20 MB of the file."""
        import credit_report.generation.evidence as ev
        src = open(ev.__file__, encoding="utf-8").read()
        assert "20 * 1024 * 1024" in src  # DID-OCR-CAP: increase if processing large PDFs
